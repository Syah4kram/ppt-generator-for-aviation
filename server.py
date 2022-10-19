import os
import sqlite3
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
from werkzeug.utils import secure_filename
from flask import Flask, render_template, request, send_file
from waitress import serve

app = Flask(__name__, static_url_path="", static_folder='static')
DATABASE = 'database.db'
UPLOAD_FOLDER = 'static/uploads/'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/adddata', methods=['GET', 'POST'])
def adddata():
    if request.method == 'POST':
        date = request.form['date']
        station = request.form['station']
        cape = request.form['cape']
        file = request.files['file']
        si = request.form['si']
        ki = request.form['ki']
        li = request.form['li']
        potensiPertumbuhanAwan = request.form['potensi-pertumbuhan-awan']
        lightning = request.form['lightning']
        hujan = request.form['hujan']
        microburst = request.form['microburst']
        uap = request.form['uap-air']
        satudari = request.form['satu-dari']
        satuhingga = request.form['satu-hingga']
        satuknots = request.form['satu-knots']
        duadari = request.form['dua-dari']
        duahingga = request.form['dua-hingga']
        duaknots = request.form['dua-knots']
        tigadari = request.form['tiga-dari']
        tigahingga = request.form['tiga-hingga']
        tigaknots = request.form['tiga-knots']
        empatdari = request.form['empat-dari']
        empathingga = request.form['empat-hingga']
        empatknots = request.form['empat-knots']
        limadari = request.form['lima-dari']
        limahingga = request.form['lima-hingga']
        limaknots = request.form['lima-knots']

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file.save(os.path.join(UPLOAD_FOLDER, filename))

            with sqlite3.connect("database.db") as con:
                cur = con.cursor()
                cur.execute(
                    "INSERT INTO data (date, station, image, cape, si, li, ki, potensiPertumbuhanAwan, lightning, hujan, microburst, uap, satudari, satuhingga, satuknots, duadari, duahingga, duaknots, tigadari, tigahingga, tigaknots, empatdari, empathingga, empatknots, limadari, limahingga, limaknots)  VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)", (date, station, filename, cape, si, li, ki, potensiPertumbuhanAwan, lightning, hujan, microburst, uap, satudari, satuhingga, satuknots, duadari, duahingga, duaknots, tigadari, tigahingga, tigaknots, empatdari, empathingga, empatknots, limadari, limahingga, limaknots))

                con.commit()
                msg = "Record successfully added"
            con.close()
            return render_template("result.html", msg=msg)

    elif request.method=='GET':
        return render_template('index.html')

    else:
        print('ERROR')

@app.route('/generateppt')
def generateppt():
    con = sqlite3.connect('database.db')

    alldata = "SELECT * FROM data"

    cur = con.cursor()
    cur.execute(alldata)
    rows = cur.fetchall()

    i = len(rows)-1

    station = rows[i][1]
    date = rows[i][0]
    img = rows[i][2]
    cape = rows[i][3]
    si = rows[i][4]
    ki = rows[i][5]
    li = rows[i][6]
    awan = rows[i][7]
    lightning = rows[i][8]
    hujan = rows[i][9]
    mb = rows[i][10]
    uap = rows[i][11]
    angin_satu_dari = rows[i][12]
    angin_satu_ke = rows[i][13]
    angin_satu_knot = rows[i][14]
    angin_dua_dari = rows[i][15]
    angin_dua_ke = rows[i][16]
    angin_dua_knot = rows[i][17]
    angin_tiga_dari = rows[i][18]
    angin_tiga_ke = rows[i][19]
    angin_tiga_knot = rows[i][20]
    angin_empat_dari = rows[i][21]
    angin_empat_ke = rows[i][22]
    angin_empat_knot = rows[i][23]
    angin_lima_dari = rows[i][24]
    angin_lima_ke = rows[i][25]
    angin_lima_knot = rows[i][26]

    con.close()

    def kondisi_udara(cape, li):
        if cape < 0 and li > 0:
            kondisi_udara = "Stabil \n (Tidak ada potensi pertumbuhan awan konvektif)"
        elif 1 <= cape <= 100 and -2 <= li <= 0:
            kondisi_udara = "Labil lemah / cukup stabil \n (Potensi pertumbuhan awan konvektif kecil)"
        elif 1000 <= cape <= 2500 and -6 <= li <= -3:
            kondisi_udara = "Labil sedang /cukup labil \n (Potensi pertumbuhan awan konvektif sedang)"
        else: kondisi_udara = "Labil kuat \n (Potensi pertumbuhan awan konvektif besar)"
        return kondisi_udara

    def thunderstorm(si, ki):
        if ki < 15:
            ts = "Tidak ada potensi TS"
        elif si < 0 and 15 <= ki <= 25:
            ts = "Potensi TS lemah"
        elif -3 <= si <= 0 and 26 <= ki <= 25:
            ts = "Potensi TS sedang"
        elif -6 <= si <= -3 and 31 <= ki <= 40:
            ts = "TS kuat"
        else: ts = "Potensi TS Sangat Kuat"
        return ts

    def judulSlide():
        #blank ppt file
        ppt_template = 'ppt.pptx'
        prs = Presentation(ppt_template)

        #slide template
        SLD_TITLE = 0
        SLD_TITLE_AND_CONTENT = 1

        slide_title = "ANALISIS DATA UDARA ATAS"
        slide_subtitle = "Pusat Meteorologi Penerbangan"
        #Add a title slide
        slide_layout = prs.slide_layouts[SLD_TITLE]
        slide = prs.slides.add_slide(slide_layout)

        #Insert title slide text. There are two text shapes on the Title slide layout by default.
        title1 = slide.shapes.title
        plc1 = slide.placeholders[1]
        title1.text = slide_title
        plc1.text = slide_subtitle

        if lightning == 0:
            light = "ada"
        else:
            light = "tidak ada"

        if hujan == 0:
            rain = "ada"
        else: rain = "tidak ada"

        if mb == 0:
            microburst = "ada"
        else: microburst = "tidak ada"

        slide_content_title = "{} \n TANGGAL {} JAM 00 UTC".format(station, date)
        content_kecepatan = '''Arah dan Kecepatan Angin \n  0-5000 ft : dari {} hingga {} : {} knots \n  5000-9000 ft : dari {} hingga {} : {} knots \n  9000-23000 ft : dari {} hingga {} : {} knots \n  23000-39000 ft : dari {} hingga {} : {} knots \n  39000-55000 ft : dari {} hingga {} : {} knots'''.format(angin_satu_dari, angin_satu_ke , angin_satu_knot , angin_dua_dari, angin_dua_ke, angin_dua_knot, angin_tiga_dari, angin_tiga_ke, angin_tiga_knot , angin_empat_dari , angin_empat_ke , angin_empat_knot, angin_lima_dari , angin_lima_ke , angin_lima_knot)
        slide_content = '''Potensi Thunderstorm: {} \n Terdapat potensi awan pada ketinggian {} ft \n Potensi Lightning: {} \n  Potensi Hujan: {}  \n Potensi Microburst: {} \n Kondisi udara {} \n Kondisi uap air: {} mm '''.format(thunderstorm(si, ki), awan, light, rain, microburst, kondisi_udara(cape, li), uap)

        slide_layout2 = prs.slide_layouts[SLD_TITLE_AND_CONTENT]
        slide2 = prs.slides.add_slide(slide_layout2)

        sld1_placeholder0 = slide2.placeholders[0]
        sld1_placeholder1 = slide2.placeholders[1]
        sld1_placeholder0.text = slide_content_title

        p = sld1_placeholder0.text_frame.paragraphs[0]
        p.font.size = Pt(23)
        p2 = sld1_placeholder0.text_frame.paragraphs[1]
        p2.font.size = Pt(15)

        left = Cm(2.25)
        top = Cm(2.85)
        height = Cm(11)
        width = Cm(22)

        #path = 'chart.jpeg'
        path = os.path.join(UPLOAD_FOLDER, img)
        slide2.shapes.add_picture(path, left, top, width, height)

        l = Cm(13.45)
        t = Cm(13.50)
        h = Cm(5.64)
        w = Cm(11.55)

        left = Cm(1)
        width = Cm(12.15)
        height = Cm(5.12)
        top = Cm(13.50)

        tb = slide2.shapes.add_textbox(l, t, w, h)
        tb.text_frame.paragraphs[0].add_run().font.size = Cm(10)
        tb.text = content_kecepatan
        tb.line.fill.solid()

        tb = slide2.shapes.add_textbox(left, top, width, height)
        tb.text = slide_content
        tb.line.fill.solid()

        prs.save('sample.pptx')

    judulSlide()
    #addSlide()
    #download()
    filepath = "sample.pptx"
    return send_file(filepath)
    #return send_file(filepath, attachment_filename = 'sample.pptx', as_attachment=True)

if __name__ == '__main__':
    #app.run(debug=True)
    serve(app, host="0.0.0.0", port=8080)
